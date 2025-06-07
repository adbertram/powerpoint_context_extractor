#!/usr/bin/env python3
"""
PowerPoint Extractor
-------------------
This script extracts slides, animations, and notes from a PowerPoint file and saves them to JSON files.
It combines the functionality of extract_slides.py, extract_animations.py, and extract_notes.py.
"""

import os
import sys
import json
import argparse
import logging
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Define XML namespaces used in PPTX files
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006'
}

def register_namespaces():
    """Register XML namespaces for parsing."""
    for prefix, uri in NAMESPACES.items():
        ET.register_namespace(prefix, uri)

def extract_notes_from_xml(pptx_path):
    """
    Extract notes directly from the PPTX XML structure.
    
    This function parses the notes XML files in a PowerPoint presentation to extract
    notes text. It looks for text within the body placeholder in notes slides.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        dict: Dictionary mapping slide numbers to notes text
    """
    notes_by_slide = {}
    
    try:
        with zipfile.ZipFile(pptx_path) as pptx_zip:
            # Get list of all notes files
            notes_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/notesSlides/notesSlide') and f.endswith('.xml')]
            
            if not notes_files:
                logger.info("No notes files found in the PowerPoint file.")
                return notes_by_slide
            
            logger.info(f"Found {len(notes_files)} notes files in the PowerPoint file.")
            
            # Process each notes file
            for notes_file in notes_files:
                try:
                    # Extract slide number from notes file name
                    slide_num = int(notes_file.split('notesSlide')[1].split('.')[0])
                    
                    # Extract notes text from XML
                    with pptx_zip.open(notes_file) as notes_xml:
                        xml_content = notes_xml.read().decode('utf-8')
                        logger.debug(f"Processing notes file: {notes_file}")
                        
                        # Parse XML content
                        root = ET.fromstring(xml_content)
                        
                        # Find all shapes in the notes slide
                        shapes = root.findall('.//p:sp', NAMESPACES)
                        
                        # Look for the shape with the notes content (has placeholder type="body")
                        notes_text = ""
                        for shape in shapes:
                            # Check if this is the notes placeholder
                            ph_elem = shape.find('.//p:nvPr/p:ph[@type="body"]', NAMESPACES)
                            if ph_elem is not None:
                                # This is the notes placeholder, extract text
                                tx_body = shape.find('.//p:txBody', NAMESPACES)
                                if tx_body is not None:
                                    # Extract all paragraphs
                                    paragraphs = []
                                    
                                    # Find all paragraph elements
                                    p_elems = tx_body.findall('.//a:p', NAMESPACES)
                                    for p in p_elems:
                                        # Find all text runs in this paragraph
                                        r_elems = p.findall('.//a:r', NAMESPACES)
                                        para_text = []
                                        
                                        # If there are no text runs, check for direct text elements
                                        if not r_elems:
                                            t_elems = p.findall('.//a:t', NAMESPACES)
                                            for t in t_elems:
                                                if t.text and t.text.strip():
                                                    para_text.append(t.text)
                                        else:
                                            # Extract text from each run
                                            for r in r_elems:
                                                t_elems = r.findall('.//a:t', NAMESPACES)
                                                for t in t_elems:
                                                    if t.text and t.text.strip():
                                                        para_text.append(t.text)
                                        
                                        # Add paragraph text if not empty
                                        if para_text:
                                            paragraphs.append(' '.join(para_text))
                                    
                                    # Combine paragraphs into notes text
                                    if paragraphs:
                                        notes_text = '\n'.join(paragraphs)
                                        notes_by_slide[slide_num] = notes_text
                                        logger.debug(f"Found notes for slide {slide_num}: {notes_text[:50]}..." if len(notes_text) > 50 else f"Found notes for slide {slide_num}: {notes_text}")
                                        break  # Found the notes, no need to check other shapes
                        
                        # Log if no notes content was found
                        if slide_num not in notes_by_slide:
                            logger.debug(f"No notes content found in {notes_file}")
                
                except Exception as e:
                    logger.error(f"Error processing notes file {notes_file}: {e}", exc_info=True)
    
    except Exception as e:
        logger.error(f"Error extracting notes from XML: {e}", exc_info=True)
    
    return notes_by_slide

def extract_slide_notes(pptx_path):
    """
    Extract notes from all slides in a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
    
    Returns:
        dict: Dictionary containing notes for all slides
    """
    # Register XML namespaces
    register_namespaces()
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Extract notes using direct XML parsing
    logger.info("Extracting notes from XML structure...")
    xml_notes = extract_notes_from_xml(pptx_path)
    
    # Dictionary to store notes information
    notes_data = {}
    
    # Process each slide
    logger.info(f"Found {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    title = shape.text.strip().replace('\n', ' ')
                    break
        
        # Extract notes text using python-pptx
        pptx_notes = ""
        if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
            if hasattr(slide.notes_slide, 'notes_text_frame'):
                pptx_notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Get notes from XML parsing
        xml_note = xml_notes.get(i, "")
        
        # Use the best available notes (prefer XML parsing if it found notes)
        notes_text = xml_note if xml_note else pptx_notes
        
        logger.info(f"Processing slide {i}: {title[:50]}{'...' if len(title) > 50 else ''} - Notes: {'Yes' if notes_text else 'No'}")
        
        # Add slide information to the dictionary
        notes_data[f"slide_{i}"] = {
            'slide_number': i,
            'title': title,
            'notes': notes_text
        }
    
    return notes_data

def extract_animation_info(slide_part):
    """
    Extract animation information from a slide's XML part.
    
    Args:
        slide_part: The slide part object from python-pptx
        
    Returns:
        list: List of dictionaries containing animation information
    """
    animations = []
    
    try:
        # Access the slide's XML through the xmlpart property
        if not hasattr(slide_part, 'xmlpart') or not hasattr(slide_part.xmlpart, 'element'):
            logger.debug(f"Slide part does not have xmlpart or element attribute")
            return animations
            
        slide_xml = slide_part.xmlpart.element
        
        # Find timing information
        timing_node = slide_xml.find('.//p:timing', NAMESPACES)
        if timing_node is None:
            return animations
        
        # Find animation sequences
        tn_lt = timing_node.find('.//p:tnLst', NAMESPACES)
        if tn_lt is None:
            return animations
        
        # Process each animation sequence
        for i, par in enumerate(tn_lt.findall('.//p:par', NAMESPACES)):
            ctn = par.find('.//p:cTn', NAMESPACES)
            if ctn is None:
                continue
                
            # Get sequence ID and duration
            seq_id = ctn.get('id', f'unknown_{i}')
            dur = ctn.get('dur', 'unknown')
            
            # Find child animations
            child_tn_lt = ctn.find('.//p:childTnLst', NAMESPACES)
            if child_tn_lt is None:
                continue
                
            # Process each animation effect
            for j, child_par in enumerate(child_tn_lt.findall('.//p:par', NAMESPACES)):
                child_ctn = child_par.find('.//p:cTn', NAMESPACES)
                if child_ctn is None:
                    continue
                    
                # Get effect ID and duration
                effect_id = child_ctn.get('id', f'unknown_effect_{j}')
                effect_dur = child_ctn.get('dur', 'unknown')
                
                # Find target shape
                tgt_el = child_par.find('.//p:tgtEl', NAMESPACES)
                if tgt_el is None:
                    continue
                    
                # Get shape ID
                shape_id_el = tgt_el.find('.//p:spTgt', NAMESPACES)
                shape_id = "unknown"
                if shape_id_el is not None:
                    shape_id = shape_id_el.get('spid', 'unknown')
                
                # Find animation effect
                anim_effect = child_par.find('.//p:animEffect', NAMESPACES)
                effect_type = "unknown"
                if anim_effect is not None:
                    effect_type = anim_effect.get('transition', anim_effect.get('type', 'unknown'))
                
                # Find start conditions
                cond = child_par.find('.//p:cond', NAMESPACES)
                trigger = "unknown"
                delay = "0"
                if cond is not None:
                    trigger = cond.get('evt', 'unknown')
                    delay = cond.get('delay', '0')
                
                # Add animation to list
                animations.append({
                    'sequence_id': seq_id,
                    'effect_id': effect_id,
                    'shape_id': shape_id,
                    'effect_type': effect_type,
                    'trigger': trigger,
                    'delay': delay,
                    'duration': effect_dur
                })
    
    except Exception as e:
        logger.error(f"Error extracting animation info: {e}", exc_info=True)
    
    return animations

def extract_slide_animations(pptx_path):
    """
    Extract animations from all slides in a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        dict: Dictionary containing animation information for all slides
    """
    # Register XML namespaces
    register_namespaces()
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file for animation extraction: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Dictionary to store animation data
    animation_data = {}
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    title = shape.text.strip().replace('\n', ' ')
                    break
        
        # Extract animations
        animations = extract_animation_info(slide.part)
        
        # Get shape information
        shape_info = {}
        for shape in slide.shapes:
            if shape.shape_id:
                shape_type = "Unknown"
                if hasattr(shape, "shape_type"):
                    shape_type = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
                
                shape_text = ""
                if hasattr(shape, "text") and shape.has_text_frame:
                    shape_text = shape.text.strip()
                
                shape_info[str(shape.shape_id)] = {
                    'type': shape_type,
                    'text': shape_text[:100] + ('...' if len(shape_text) > 100 else '')
                }
        
        # Get slide transition
        transition = "None"
        if hasattr(slide, "slide_layout") and hasattr(slide.slide_layout, "transition"):
            transition = str(slide.slide_layout.transition)
        
        # Add slide information to the dictionary
        animation_data[f"slide_{i}"] = {
            'slide_number': i,
            'title': title,
            'animations': animations,
            'shapes': shape_info,
            'transition': transition,
            'animation_count': len(animations)
        }
        
        logger.info(f"Processed slide {i}: {title[:50]}{'...' if len(title) > 50 else ''} - Animations: {len(animations)}")
    
    return animation_data

def extract_pptx_content(pptx_path, output_dir):
    """
    Extract slides, animations, and notes from a PowerPoint file and save to JSON files.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save output files
        
    Returns:
        tuple: (notes_data, animation_data) dictionaries
    """
    # Create output directory if it doesn't exist
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Extract notes
    logger.info("Extracting slide notes...")
    notes_data = extract_slide_notes(pptx_path)
    
    # Extract animations
    logger.info("Extracting slide animations...")
    animation_data = extract_slide_animations(pptx_path)
    
    # Save notes information to JSON file
    notes_file = output_path / "slide_notes.json"
    logger.info(f"Saving notes information to {notes_file}")
    try:
        with open(notes_file, 'w', encoding='utf-8') as f:
            json.dump(notes_data, f, indent=2)
        logger.info(f"Successfully saved notes information to {notes_file}")
    except Exception as e:
        logger.error(f"Failed to save notes information: {e}")
    
    # Save animation information to JSON file
    animation_file = output_path / "slide_animations.json"
    logger.info(f"Saving animation information to {animation_file}")
    try:
        with open(animation_file, 'w', encoding='utf-8') as f:
            json.dump(animation_data, f, indent=2)
        logger.info(f"Successfully saved animation information to {animation_file}")
    except Exception as e:
        logger.error(f"Failed to save animation information: {e}")
    
    # Create a combined JSON with both notes and animations
    combined_data = {}
    for slide_key in notes_data:
        combined_data[slide_key] = {
            'slide_number': notes_data[slide_key]['slide_number'],
            'title': notes_data[slide_key]['title'],
            'notes': notes_data[slide_key]['notes'],
            'animations': animation_data.get(slide_key, {}).get('animations', []),
            'animation_count': animation_data.get(slide_key, {}).get('animation_count', 0)
        }
    
    # Save combined information to JSON file
    combined_file = output_path / "slide_content.json"
    logger.info(f"Saving combined information to {combined_file}")
    try:
        with open(combined_file, 'w', encoding='utf-8') as f:
            json.dump(combined_data, f, indent=2)
        logger.info(f"Successfully saved combined information to {combined_file}")
    except Exception as e:
        logger.error(f"Failed to save combined information: {e}")
    
    return notes_data, animation_data

def main():
    parser = argparse.ArgumentParser(description='Extract slides, animations, and notes from a PowerPoint file.')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    parser.add_argument('--output', '-o', default='./output', help='Output directory')
    
    args = parser.parse_args()
    
    # Check if the PowerPoint file exists
    if not os.path.isfile(args.pptx_file):
        logger.error(f"PowerPoint file not found: {args.pptx_file}")
        sys.exit(1)
    
    # Extract PowerPoint content
    extract_pptx_content(args.pptx_file, args.output)

if __name__ == '__main__':
    main()
