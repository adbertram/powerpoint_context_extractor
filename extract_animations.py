#!/usr/bin/env python3
"""
PowerPoint Animation Extractor
-----------------------------
This script extracts animation information from a PowerPoint file and saves it to a JSON file.
"""

import os
import sys
import json
import argparse
import logging
from pathlib import Path
from pptx import Presentation
import xml.etree.ElementTree as ET

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Define XML namespaces used in PPTX files
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006'
}

def register_namespaces():
    """Register XML namespaces for parsing."""
    for prefix, uri in NAMESPACES.items():
        ET.register_namespace(prefix, uri)

def extract_animation_info(slide_part):
    """
    Extract animation information from a slide.
    
    Args:
        slide_part: The slide part object from python-pptx
        
    Returns:
        dict: Dictionary containing animation information
    """
    animations = []
    
    # Check if the slide has timing information
    if not hasattr(slide_part, 'element') or not hasattr(slide_part.element, 'get_or_add_timing'):
        return animations
    
    try:
        # Access the XML directly to get animation data
        xml_bytes = slide_part.blob
        root = ET.fromstring(xml_bytes)
        
        # Look for timing information
        timing_nodes = root.findall('.//p:timing', NAMESPACES)
        if not timing_nodes:
            return animations
        
        # Extract animation sequences
        for timing in timing_nodes:
            # Find animation sequences
            anim_seq_nodes = timing.findall('.//p:seq', NAMESPACES)
            
            for seq_idx, seq in enumerate(anim_seq_nodes):
                # Get animation nodes
                anim_nodes = seq.findall('.//p:anim', NAMESPACES) + seq.findall('.//p:animEffect', NAMESPACES)
                
                for anim_idx, anim in enumerate(anim_nodes):
                    animation = {
                        'sequence': seq_idx + 1,
                        'order': anim_idx + 1,
                        'type': 'unknown'
                    }
                    
                    # Try to get animation type
                    if 'type' in anim.attrib:
                        animation['type'] = anim.attrib['type']
                    
                    # Try to get animation effect
                    effect = anim.get('{{{0}}}effect'.format(NAMESPACES['p']))
                    if effect:
                        animation['effect'] = effect
                    
                    # Try to get target shape
                    target_nodes = anim.findall('.//p:tgtEl', NAMESPACES)
                    if target_nodes:
                        for target in target_nodes:
                            shape_nodes = target.findall('.//p:spTgt', NAMESPACES)
                            if shape_nodes and '{{{0}}}spid'.format(NAMESPACES['p']) in shape_nodes[0].attrib:
                                animation['target_shape_id'] = shape_nodes[0].attrib['{{{0}}}spid'.format(NAMESPACES['p'])]
                    
                    # Try to get timing information
                    cond_nodes = anim.findall('./p:cond', NAMESPACES) or anim.findall('../p:cond', NAMESPACES)
                    if cond_nodes:
                        for cond in cond_nodes:
                            if 'delay' in cond.attrib:
                                animation['delay'] = cond.attrib['delay']
                            if 'evt' in cond.attrib:
                                animation['trigger_event'] = cond.attrib['evt']
                    
                    animations.append(animation)
                
                # Look for p:par nodes which can contain animation information
                par_nodes = seq.findall('.//p:par', NAMESPACES)
                for par_idx, par in enumerate(par_nodes):
                    ctn_nodes = par.findall('.//p:cTn', NAMESPACES)
                    for ctn in ctn_nodes:
                        if 'id' in ctn.attrib:
                            # Look for child animations
                            child_anims = par.findall('.//p:anim', NAMESPACES) + par.findall('.//p:animEffect', NAMESPACES)
                            for child_idx, child_anim in enumerate(child_anims):
                                animation = {
                                    'sequence': seq_idx + 1,
                                    'par_group': par_idx + 1,
                                    'order': child_idx + 1,
                                    'type': 'unknown'
                                }
                                
                                # Try to get animation type
                                if 'type' in child_anim.attrib:
                                    animation['type'] = child_anim.attrib['type']
                                
                                # Try to get animation effect
                                effect = child_anim.get('{{{0}}}effect'.format(NAMESPACES['p']))
                                if effect:
                                    animation['effect'] = effect
                                
                                animations.append(animation)
        
        # Look for transition information
        transition_nodes = root.findall('.//p:transition', NAMESPACES)
        if transition_nodes:
            for transition in transition_nodes:
                transition_info = {
                    'type': 'transition'
                }
                
                # Get transition type
                for key in transition.attrib:
                    if key != '{{{0}}}spd'.format(NAMESPACES['p']):
                        transition_info['transition_type'] = key
                
                # Get transition speed
                if '{{{0}}}spd'.format(NAMESPACES['p']) in transition.attrib:
                    transition_info['speed'] = transition.attrib['{{{0}}}spd'.format(NAMESPACES['p'])]
                
                animations.append(transition_info)
    
    except Exception as e:
        logger.error(f"Error extracting animation info: {e}")
    
    return animations

def extract_slide_animations(pptx_path, output_file):
    """
    Extract animation information from all slides in a PowerPoint file and save to JSON.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_file (str): Path to save the JSON output
    
    Returns:
        dict: Dictionary containing animation information for all slides
    """
    # Register XML namespaces
    register_namespaces()
    
    # Create output directory if it doesn't exist
    output_path = Path(output_file).parent
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Dictionary to store animation information
    animation_data = {}
    
    # Process each slide
    logger.info(f"Found {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    title = shape.text.strip().replace('\n', ' ')
                    break
        
        logger.info(f"Processing slide {i}: {title}")
        
        # Extract animation information
        animations = extract_animation_info(slide.part)
        
        # Add shape information to help identify animation targets
        shape_info = []
        for shape_idx, shape in enumerate(slide.shapes):
            shape_data = {
                'id': shape.shape_id,
                'name': shape.name,
                'type': shape.shape_type
            }
            if hasattr(shape, "text") and shape.has_text_frame:
                shape_data['text'] = shape.text.strip()
            shape_info.append(shape_data)
        
        # Add slide information to the dictionary
        animation_data[f"slide_{i}"] = {
            'title': title,
            'animations': animations,
            'shapes': shape_info
        }
    
    # Save animation information to JSON file
    logger.info(f"Saving animation information to {output_file}")
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(animation_data, f, indent=2)
        logger.info(f"Successfully saved animation information to {output_file}")
    except Exception as e:
        logger.error(f"Failed to save animation information: {e}")
        return None
    
    return animation_data

def main():
    parser = argparse.ArgumentParser(description='Extract animation information from a PowerPoint file.')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    parser.add_argument('--output', '-o', default='./animations.json', help='Output JSON file')
    
    args = parser.parse_args()
    
    # Check if the PowerPoint file exists
    if not os.path.isfile(args.pptx_file):
        logger.error(f"PowerPoint file not found: {args.pptx_file}")
        sys.exit(1)
    
    result = extract_slide_animations(args.pptx_file, args.output)
    if result is None:
        sys.exit(1)

if __name__ == '__main__':
    main()
