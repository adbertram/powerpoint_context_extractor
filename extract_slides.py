#!/usr/bin/env python3
"""
PowerPoint Slide Extractor
--------------------------
This script extracts slides from a PowerPoint file and saves each slide as an image.
The image filenames are formatted as "slide_number-slide_title.png".
"""

import os
import re
import sys
import argparse
import subprocess
import tempfile
import shutil
from pathlib import Path
from pptx import Presentation
import logging
from pdf2image import convert_from_path

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def sanitize_filename(filename):
    """Sanitize the filename by removing invalid characters."""
    # Replace invalid filename characters with underscores
    sanitized = re.sub(r'[\\/*?:"<>|]', '_', filename)
    # Replace multiple spaces with a single underscore
    sanitized = re.sub(r'\s+', '_', sanitized)
    # Limit filename length
    if len(sanitized) > 100:
        sanitized = sanitized[:97] + '...'
    return sanitized

def check_dependencies():
    """Check if required external dependencies are installed."""
    dependencies = {
        'soffice': 'LibreOffice is not installed. Install it with: brew install libreoffice',
        'unoconv': 'unoconv is not installed. Install it with: brew install unoconv',
        'pdftoppm': 'poppler is not installed. Install it with: brew install poppler'
    }
    
    missing = []
    for cmd, message in dependencies.items():
        try:
            subprocess.run(['which', cmd], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        except subprocess.CalledProcessError:
            missing.append(message)
    
    return missing

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """
    Convert PowerPoint file to PDF using LibreOffice or unoconv.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        pdf_path (str): Path to save the PDF file
    
    Returns:
        bool: True if conversion was successful, False otherwise
    """
    try:
        # First try using LibreOffice
        logger.info(f"Converting {pptx_path} to PDF using LibreOffice...")
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", os.path.dirname(pdf_path),
            pptx_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode == 0:
            # Rename the output file if needed
            default_pdf = os.path.join(
                os.path.dirname(pdf_path),
                os.path.basename(pptx_path).rsplit('.', 1)[0] + '.pdf'
            )
            if default_pdf != pdf_path:
                shutil.move(default_pdf, pdf_path)
            return True
        else:
            logger.warning(f"LibreOffice conversion failed: {result.stderr}")
            
            # Try using unoconv as fallback
            logger.info("Trying unoconv as fallback...")
            cmd = ["unoconv", "-f", "pdf", "-o", pdf_path, pptx_path]
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                return True
            else:
                logger.error(f"unoconv conversion failed: {result.stderr}")
                return False
    
    except Exception as e:
        logger.error(f"PDF conversion error: {e}")
        return False

def extract_slide_titles(pptx_path):
    """
    Extract slide titles from PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
    
    Returns:
        list: List of slide titles
    """
    titles = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            title = "Untitled"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    if shape.text.strip():
                        title = shape.text.strip().replace('\n', ' ')
                        break
            titles.append(title)
        return titles
    except Exception as e:
        logger.error(f"Failed to extract slide titles: {e}")
        return []

def extract_slides(pptx_path, output_dir, format='png', dpi=300):
    """
    Extract slides from a PowerPoint file and save them as images.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save the images
        format (str): Image format (default: png)
        dpi (int): Image resolution in DPI (default: 300)
    
    Returns:
        Path: Path to the output directory
    """
    # Check dependencies
    missing_deps = check_dependencies()
    if missing_deps:
        logger.error("Missing dependencies:")
        for dep in missing_deps:
            logger.error(f"  - {dep}")
        logger.error("Please install the missing dependencies and try again.")
        return None
    
    # Create output directory if it doesn't exist
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Extract slide titles
    logger.info("Extracting slide titles...")
    titles = extract_slide_titles(pptx_path)
    
    if not titles:
        logger.error("Failed to extract slide titles")
        return None
    
    logger.info(f"Found {len(titles)} slides")
    
    # Create temporary directory for PDF
    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "presentation.pdf")
        
        # Convert PowerPoint to PDF
        if not convert_pptx_to_pdf(pptx_path, pdf_path):
            logger.error("Failed to convert PowerPoint to PDF")
            return None
        
        # Convert PDF to images
        logger.info("Converting PDF to images...")
        try:
            images = convert_from_path(pdf_path, dpi=dpi)
            
            # Save each image with the corresponding slide title
            for i, (image, title) in enumerate(zip(images, titles), 1):
                safe_title = sanitize_filename(title)
                filename = f"{i:02d}-{safe_title}.{format}"
                output_file = output_path / filename
                
                logger.info(f"Saving slide {i}: {title}")
                image.save(str(output_file), format.upper())
                
            logger.info(f"Successfully extracted {len(images)} slides to {output_path}")
            
        except Exception as e:
            logger.error(f"Failed to convert PDF to images: {e}")
            return None
    
    return output_path

def main():
    parser = argparse.ArgumentParser(description='Extract slides from a PowerPoint file as images.')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    parser.add_argument('--output', '-o', default='./slides', help='Output directory for images')
    parser.add_argument('--format', '-f', default='png', choices=['png', 'jpg', 'jpeg'], help='Image format')
    parser.add_argument('--dpi', '-d', type=int, default=300, help='Image resolution in DPI')
    
    args = parser.parse_args()
    
    # Check if the PowerPoint file exists
    if not os.path.isfile(args.pptx_file):
        logger.error(f"PowerPoint file not found: {args.pptx_file}")
        sys.exit(1)
    
    result = extract_slides(args.pptx_file, args.output, args.format, args.dpi)
    if result is None:
        sys.exit(1)

if __name__ == '__main__':
    main()
