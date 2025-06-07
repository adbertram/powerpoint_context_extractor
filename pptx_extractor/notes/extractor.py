"""
Notes extraction functionality for PowerPoint presentations.
"""

import re
import zipfile
import xml.etree.ElementTree as ET
import logging
from pptx import Presentation

from ..utils.common import NAMESPACES, register_namespaces

logger = logging.getLogger(__name__)

def extract_notes_from_xml(pptx_path):
    """
    Extract notes directly from the PPTX XML structure.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        dict: Dictionary mapping slide numbers to notes text
    """
    notes_by_slide = {}
    
    try:
        with zipfile.ZipFile(pptx_path) as pptx_zip:
            # Get list of all notes files
            notes_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/notesSlides/notesSlide') and f.endswith('.xml')]
            
            if not notes_files:
                logger.info("No notes files found in the PowerPoint file.")
                return notes_by_slide
            
            logger.info(f"Found {len(notes_files)} notes files in the PowerPoint file.")
            
            # Process each notes file
            for notes_file in notes_files:
                try:
                    # Extract slide number from notes file name
                    import re
                    match = re.search(r'notesSlide(\d+)\.xml', notes_file)
                    if not match:
                        continue
                    slide_num = int(match.group(1))
                    
                    # Extract notes text from XML
                    with pptx_zip.open(notes_file) as notes_xml:
                        xml_content = notes_xml.read().decode('utf-8')
                        logger.debug(f"Processing notes file: {notes_file}")
                        
                        # Parse XML content
                        root = ET.fromstring(xml_content)
                        
                        # Find all shapes in the notes slide
                        shapes = root.findall('.//p:sp', NAMESPACES)
                        
                        # Look for the shape with the notes content (has placeholder type="body")
                        notes_text = ""
                        for shape in shapes:
                            # Check if this is the notes placeholder
                            ph_elem = shape.find('.//p:nvPr/p:ph[@type="body"]', NAMESPACES)
                            if ph_elem is not None:
                                # This is the notes placeholder, extract text
                                tx_body = shape.find('.//p:txBody', NAMESPACES)
                                if tx_body is not None:
                                    # Extract all paragraphs
                                    paragraphs = []
                                    
                                    # Find all paragraph elements
                                    p_elems = tx_body.findall('.//a:p', NAMESPACES)
                                    for p in p_elems:
                                        # Find all text runs in this paragraph
                                        r_elems = p.findall('.//a:r', NAMESPACES)
                                        para_text = []
                                        
                                        # If there are no text runs, check for direct text elements
                                        if not r_elems:
                                            t_elems = p.findall('.//a:t', NAMESPACES)
                                            for t in t_elems:
                                                if t.text and t.text.strip():
                                                    para_text.append(t.text)
                                        else:
                                            # Extract text from each run
                                            for r in r_elems:
                                                t_elems = r.findall('.//a:t', NAMESPACES)
                                                for t in t_elems:
                                                    if t.text and t.text.strip():
                                                        para_text.append(t.text)
                                        
                                        # Add paragraph text if not empty
                                        if para_text:
                                            paragraphs.append(' '.join(para_text))
                                    
                                    # Combine paragraphs into notes text
                                    if paragraphs:
                                        notes_text = '\n'.join(paragraphs)
                                        notes_by_slide[slide_num] = notes_text
                                        logger.debug(f"Found notes for slide {slide_num}: {notes_text[:50]}..." if len(notes_text) > 50 else f"Found notes for slide {slide_num}: {notes_text}")
                                        break  # Found the notes, no need to check other shapes
                        
                        # Log if no notes content was found
                        if slide_num not in notes_by_slide:
                            logger.debug(f"No notes content found in {notes_file}")
                
                except Exception as e:
                    logger.error(f"Error processing notes file {notes_file}: {e}", exc_info=True)
    
    except Exception as e:
        logger.error(f"Error extracting notes from XML: {e}", exc_info=True)
    
    return notes_by_slide

def extract_slide_notes(pptx_path):
    """
    Extract notes from all slides in a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
    
    Returns:
        dict: Dictionary containing notes for all slides
    """
    # Register XML namespaces
    register_namespaces()
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Extract notes using direct XML parsing
    logger.info("Extracting notes from XML structure...")
    xml_notes = extract_notes_from_xml(pptx_path)
    
    # Dictionary to store notes information
    notes_data = {}
    
    # Process each slide
    logger.info(f"Found {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    title = shape.text.strip().replace('\n', ' ')
                    break
        
        # Extract notes text using python-pptx
        pptx_notes = ""
        if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
            if hasattr(slide.notes_slide, 'notes_text_frame'):
                pptx_notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Get notes from XML parsing
        xml_note = xml_notes.get(i, "")
        
        # Use the best available notes (prefer XML parsing if it found notes)
        notes_text = xml_note if xml_note else pptx_notes
        
        logger.info(f"Processing slide {i}: {title[:50]}{'...' if len(title) > 50 else ''} - Notes: {'Yes' if notes_text else 'No'}")
        
        # Add slide information to the dictionary
        notes_data[f"slide_{i}"] = {
            'slide_number': i,
            'title': title,
            'notes': notes_text
        }
    
    return notes_data
