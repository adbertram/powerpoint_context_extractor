"""
Slides extraction functionality for PowerPoint presentations.
"""

import os
import subprocess
import logging
import shutil
import tempfile
from pathlib import Path
from pptx import Presentation

from ..utils.common import ensure_directory, sanitize_filename, get_slide_title

logger = logging.getLogger(__name__)

def check_dependencies():
    """Check if required dependencies are installed.
    
    Returns:
        bool: True if all dependencies are installed, False otherwise
    """
    try:
        # Check if LibreOffice is installed
        subprocess.run(['which', 'soffice'], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Check if Poppler is installed (for pdf2image)
        subprocess.run(['which', 'pdftoppm'], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        return True
    except subprocess.CalledProcessError:
        logger.error("Required dependencies not found. Please install LibreOffice and Poppler.")
        return False
    except Exception as e:
        logger.error(f"Error checking dependencies: {e}")
        return False

def convert_pptx_to_pdf(pptx_path, output_dir):
    """Convert PowerPoint file to PDF using LibreOffice.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save the PDF file
        
    Returns:
        str: Path to the PDF file or None if conversion failed
    """
    try:
        import time
        from pptx import Presentation
        
        # Get slide count for progress reporting
        try:
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            logger.info(f"Starting conversion of PowerPoint with {slide_count} slides to PDF")
        except Exception as e:
            logger.warning(f"Could not determine slide count: {e}")
            slide_count = "unknown number of"
        
        logger.info(f"Converting {slide_count} slides from PowerPoint to PDF...")
        logger.info("This may take some time for large presentations")
        
        start_time = time.time()
        
        # Create a temporary directory for the conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert PPTX to PDF
            logger.info("Launching LibreOffice for conversion (this is a background process)...")
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ]
            
            # Run the conversion process with progress updates
            process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            
            # Provide periodic updates while waiting for conversion
            while process.poll() is None:
                logger.info("PDF conversion in progress... (this may take several minutes)")
                time.sleep(5)  # Check status every 5 seconds
            
            # Check if conversion was successful
            if process.returncode != 0:
                stderr = process.stderr.read().decode('utf-8', errors='ignore')
                logger.error(f"LibreOffice conversion failed with error: {stderr}")
                return None
            
            # Get the PDF file name
            pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
            pdf_path = os.path.join(temp_dir, pdf_filename)
            
            if not os.path.exists(pdf_path):
                logger.error(f"PDF file not created at expected path: {pdf_path}")
                return None
            
            # Copy the PDF file to the output directory
            output_pdf = os.path.join(output_dir, pdf_filename)
            shutil.copy2(pdf_path, output_pdf)
            
            conversion_time = time.time() - start_time
            logger.info(f"Successfully converted {pptx_path} to {output_pdf} in {conversion_time:.2f} seconds")
            return output_pdf
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        return None

def extract_slide_titles(pptx_path):
    """Extract titles from all slides in a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        list: List of slide titles
    """
    titles = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            title = get_slide_title(slide)
            titles.append(title)
        logger.info(f"Extracted {len(titles)} slide titles")
        return titles
    except Exception as e:
        logger.error(f"Error extracting slide titles: {e}")
        return []

def convert_pdf_to_images(pdf_path, output_dir, format='png', dpi=300):
    """Convert PDF file to images using pdf2image.
    
    Args:
        pdf_path (str): Path to the PDF file
        output_dir (str): Directory to save the images
        format (str): Image format (default: png)
        dpi (int): Image resolution (default: 300)
        
    Returns:
        list: List of paths to the generated images
    """
    try:
        import time
        import threading
        import sys
        
        logger.info(f"Starting PDF to image conversion for {pdf_path}")
        logger.info(f"This process may take several minutes for large presentations...")
        
        # Get the number of pages in the PDF (to provide progress updates)
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            doc.close()
            logger.info(f"Found {total_pages} pages in the PDF to convert")
        except ImportError:
            logger.info("PyMuPDF not installed, cannot determine page count in advance")
            total_pages = None
        except Exception as e:
            logger.warning(f"Could not determine PDF page count: {e}")
            total_pages = None
        
        # Set up a progress reporting thread
        stop_progress_thread = threading.Event()
        
        def progress_reporter():
            progress_counter = 0
            while not stop_progress_thread.is_set():
                progress_counter += 1
                logger.info(f"PDF to image conversion in progress... ({progress_counter * 5} seconds elapsed)")
                time.sleep(5)
        
        # Start progress reporting thread
        progress_thread = threading.Thread(target=progress_reporter)
        progress_thread.daemon = True
        progress_thread.start()
        
        # Convert PDF to images with progress updates
        start_time = time.time()
        logger.info("Beginning image conversion (this is CPU-intensive and may take time)...")
        
        try:
            # Import here to avoid slowing down the script if not needed
            from pdf2image import convert_from_path
            
            # Convert PDF to images
            images = convert_from_path(
                pdf_path, 
                dpi=dpi,
                # Use a thread-safe callback to report progress during conversion
                thread_count=1  # Use single thread to avoid potential issues
            )
            
            # Stop progress reporting
            stop_progress_thread.set()
            progress_thread.join(timeout=1.0)
            
            conversion_time = time.time() - start_time
            logger.info(f"PDF conversion completed in {conversion_time:.2f} seconds")
            logger.info(f"Generated {len(images)} images, now saving to disk...")
            
            # Save images with progress updates
            image_paths = []
            for i, image in enumerate(images):
                if i % 5 == 0 or i == len(images) - 1:
                    logger.info(f"Saving image {i+1}/{len(images)}")
                
                image_path = os.path.join(output_dir, f"slide_{i+1}.{format}")
                image.save(image_path, format.upper())
                image_paths.append(image_path)
                
                # Provide more frequent progress updates
                if (i+1) % 10 == 0:
                    logger.info(f"Progress: {i+1}/{len(images)} images saved")
            
            total_time = time.time() - start_time
            logger.info(f"Converted {pdf_path} to {len(image_paths)} images in {total_time:.2f} seconds")
            return image_paths
            
        except Exception as e:
            # Stop progress reporting in case of error
            stop_progress_thread.set()
            if progress_thread.is_alive():
                progress_thread.join(timeout=1.0)
            raise e
            
    except ImportError:
        logger.error("pdf2image not installed. Please install it with 'pip install pdf2image'")
        return []
    except Exception as e:
        logger.error(f"Error converting PDF to images: {e}")
        return []

def extract_slides(pptx_path, output_dir, format='png', dpi=300):
    """Extract slides from a PowerPoint file as images.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save the images
        format (str): Image format (default: png)
        dpi (int): Image resolution (default: 300)
        
    Returns:
        list: List of paths to the generated images
    """
    # Create output directory if it doesn't exist
    output_path = ensure_directory(output_dir)
    
    # Check dependencies
    if not check_dependencies():
        logger.error("Required dependencies not found. Aborting slide extraction.")
        return []
    
    # Extract slide titles
    titles = extract_slide_titles(pptx_path)
    
    # Convert PowerPoint to PDF
    pdf_path = convert_pptx_to_pdf(pptx_path, output_dir)
    if not pdf_path:
        logger.error("Failed to convert PowerPoint to PDF. Aborting slide extraction.")
        return []
    
    # Convert PDF to images
    temp_images = convert_pdf_to_images(pdf_path, output_dir, format, dpi)
    if not temp_images:
        logger.error("Failed to convert PDF to images. Aborting slide extraction.")
        return []
    
    # Rename images with slide titles
    renamed_images = []
    for i, image_path in enumerate(temp_images):
        if i < len(titles):
            title = titles[i]
            sanitized_title = sanitize_filename(title)
            new_name = f"slide_{i+1:03d}-{sanitized_title}.{format}"
            new_path = os.path.join(output_dir, new_name)
            try:
                os.rename(image_path, new_path)
                renamed_images.append(new_path)
            except Exception as e:
                logger.error(f"Error renaming image {image_path}: {e}")
                renamed_images.append(image_path)
        else:
            renamed_images.append(image_path)
    
    logger.info(f"Extracted {len(renamed_images)} slides to {output_dir}")
    return renamed_images
