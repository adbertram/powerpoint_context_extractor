"""
Animations extraction functionality for PowerPoint presentations.
"""

import logging
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

from ..utils.common import NAMESPACES, register_namespaces, get_slide_title

logger = logging.getLogger(__name__)

def has_animations_in_xml(xml_element):
    """
    Check if a slide XML element contains animation definitions.
    
    Args:
        xml_element: The XML element to check
        
    Returns:
        bool: True if animations are found, False otherwise
    """
    # Check for timing information which contains animations
    timing_node = xml_element.find('.//p:timing', NAMESPACES)
    if timing_node is None:
        return False
    
    # Check for animation sequences
    tn_lt = timing_node.find('.//p:tnLst', NAMESPACES)
    if tn_lt is None or len(list(tn_lt.findall('.//p:par', NAMESPACES))) == 0:
        return False
    
    return True

def extract_animation_info(slide):
    """
    Extract animation information from a slide.
    
    Args:
        slide: The slide object from python-pptx
        
    Returns:
        list: List of dictionaries containing animation information
    """
    animations = []
    
    try:
        # Access the slide's XML element
        if hasattr(slide, 'element'):
            slide_xml = slide.element
        elif hasattr(slide, '_element'):
            slide_xml = slide._element
        else:
            logger.debug(f"Slide does not have element attribute")
            return animations
        
        # Find timing information
        timing_node = slide_xml.find('.//p:timing', NAMESPACES)
        if timing_node is None:
            return animations
        
        # Find animation sequences
        tn_lt = timing_node.find('.//p:tnLst', NAMESPACES)
        if tn_lt is None:
            return animations
        
        # Process each animation sequence
        for i, par in enumerate(tn_lt.findall('.//p:par', NAMESPACES)):
            ctn = par.find('.//p:cTn', NAMESPACES)
            if ctn is None:
                continue
                
            # Get sequence ID and duration
            seq_id = ctn.get('id', f'unknown_{i}')
            dur = ctn.get('dur', 'unknown')
            
            # Find child animations
            child_tn_lt = ctn.find('.//p:childTnLst', NAMESPACES)
            if child_tn_lt is None:
                continue
                
            # Process each animation effect
            for j, child_par in enumerate(child_tn_lt.findall('.//p:par', NAMESPACES)):
                child_ctn = child_par.find('.//p:cTn', NAMESPACES)
                if child_ctn is None:
                    continue
                    
                # Get effect ID and duration
                effect_id = child_ctn.get('id', f'unknown_effect_{j}')
                effect_dur = child_ctn.get('dur', 'unknown')
                
                # Convert duration to milliseconds if it's a number
                duration_ms = "unknown"
                if effect_dur != 'unknown' and effect_dur.isdigit():
                    duration_ms = int(effect_dur)
                
                # Find target shape
                tgt_el = child_par.find('.//p:tgtEl', NAMESPACES)
                if tgt_el is None:
                    continue
                    
                # Get shape ID and check for paragraph target
                shape_id_el = tgt_el.find('.//p:spTgt', NAMESPACES)
                shape_id = "unknown"
                build_level = None
                if shape_id_el is not None:
                    shape_id = shape_id_el.get('spid', 'unknown')
                    # Check for text animation (by paragraph)
                    txEl = shape_id_el.find('.//p:txEl', NAMESPACES)
                    if txEl is not None:
                        pRg = txEl.find('.//p:pRg', NAMESPACES)
                        if pRg is not None:
                            build_level = f"paragraph_{pRg.get('st', '0')}-{pRg.get('end', '0')}"
                
                # Find animation effect and its properties
                anim_effect = child_par.find('.//p:animEffect', NAMESPACES)
                effect_type = "appear"  # default
                effect_subtype = None
                effect_direction = None
                
                if anim_effect is not None:
                    # Get transition type (for entrance/exit effects)
                    effect_type = anim_effect.get('transition', 'in')
                    filter_attr = anim_effect.get('filter', '')
                    
                    # Parse filter for effect details
                    if filter_attr:
                        # Common patterns: "fade", "wipe(right)", "fly(fromBottom)"
                        if '(' in filter_attr:
                            effect_name = filter_attr.split('(')[0]
                            effect_params = filter_attr.split('(')[1].rstrip(')')
                            effect_subtype = effect_name
                            effect_direction = effect_params
                        else:
                            effect_subtype = filter_attr
                
                # Check for other animation types
                if not anim_effect:
                    # Check for emphasis effects (color change, etc.)
                    anim_clr = child_par.find('.//p:animClr', NAMESPACES)
                    if anim_clr:
                        effect_type = "emphasis"
                        effect_subtype = "color"
                        # Get color details if needed
                        to_clr = anim_clr.find('.//p:to', NAMESPACES)
                        if to_clr:
                            rgb = to_clr.find('.//a:srgbClr', NAMESPACES)
                            if rgb is not None:
                                effect_direction = f"to_color_{rgb.get('val', '')}"
                    
                    # Check for motion path
                    anim_motion = child_par.find('.//p:animMotion', NAMESPACES)
                    if anim_motion:
                        effect_type = "motion"
                        effect_subtype = "path"
                        path = anim_motion.get('path', '')
                        if path:
                            effect_direction = "custom_path"
                    
                    # Check for scale/rotate
                    anim_scale = child_par.find('.//p:animScale', NAMESPACES)
                    if anim_scale:
                        effect_type = "emphasis"
                        effect_subtype = "grow/shrink"
                        by_x = anim_scale.find('.//p:by', NAMESPACES)
                        if by_x is not None:
                            x_val = by_x.get('x', '100000')
                            y_val = by_x.get('y', '100000')
                            effect_direction = f"scale_x{x_val}_y{y_val}"
                
                # Find start conditions
                start_condition = "on_click"  # default
                delay_ms = 0
                
                # Check all conditions
                stCondLst = child_ctn.find('.//p:stCondLst', NAMESPACES)
                if stCondLst:
                    cond = stCondLst.find('.//p:cond', NAMESPACES)
                    if cond is not None:
                        evt = cond.get('evt', '')
                        delay = cond.get('delay', '0')
                        
                        # Parse trigger
                        if evt == 'onBegin':
                            start_condition = "with_previous"
                        elif evt == 'onClick':
                            start_condition = "on_click"
                        elif delay == 'indefinite':
                            start_condition = "on_click"
                        else:
                            # Check for "after previous" by looking at tn
                            tn = cond.find('.//p:tn', NAMESPACES)
                            if tn is not None:
                                val = tn.get('val', '')
                                if val == 'indefinite':
                                    start_condition = "after_previous"
                        
                        # Parse delay
                        if delay and delay != 'indefinite' and delay.isdigit():
                            delay_ms = int(delay)
                
                # Get node type (main sequence, trigger, etc.)
                node_type = ctn.get('nodeType', 'mainSeq')
                
                # Get repeat and other properties
                repeat_count = ctn.get('repeatCount', '1')
                auto_reverse = ctn.get('autoRev', '0') == '1'
                
                # Add animation to list
                animations.append({
                    'sequence_id': seq_id,
                    'effect_id': effect_id,
                    'shape_id': shape_id,
                    'effect_type': effect_type,
                    'effect_subtype': effect_subtype,
                    'effect_direction': effect_direction,
                    'start_condition': start_condition,
                    'delay_ms': delay_ms,
                    'duration_ms': duration_ms,
                    'build_level': build_level,
                    'node_type': node_type,
                    'repeat_count': repeat_count,
                    'auto_reverse': auto_reverse
                })
    
    except Exception as e:
        logger.error(f"Error extracting animation info: {e}", exc_info=True)
    
    return animations

def create_animation_description(animation, shape_info):
    """
    Create a comprehensive, human-readable description of an animation for LLM understanding.
    
    Args:
        animation: Animation dictionary
        shape_info: Dictionary of shape information
        
    Returns:
        str: Detailed human-readable animation description
    """
    # Get detailed shape information
    shape_id = animation['shape_id']
    shape_details = shape_info.get(shape_id, {})
    shape_type = shape_details.get('type', 'element').replace('_', ' ').lower()
    shape_text = shape_details.get('text', '')
    
    # Create element description
    element_desc = ""
    if shape_text:
        # Truncate very long text but keep it meaningful
        display_text = shape_text[:100].strip()
        if len(shape_text) > 100:
            display_text += "..."
        
        if shape_type in ['text box', 'text placeholder', 'title']:
            element_desc = f'The {shape_type} containing "{display_text}"'
        elif shape_type == 'picture':
            element_desc = f'An image/picture element'
        else:
            element_desc = f'A {shape_type} element with content "{display_text}"'
    else:
        if shape_type == 'picture':
            element_desc = 'An image/picture element'
        elif shape_type in ['auto shape', 'freeform']:
            element_desc = 'A shape element'
        else:
            element_desc = f'A {shape_type} element'
    
    # Describe the animation effect in natural language
    effect_descriptions = {
        'in': {
            'fade': 'fades into view',
            'fly': 'flies in',
            'wipe': 'wipes in',
            'zoom': 'zooms in',
            'swivel': 'swivels in',
            'bounce': 'bounces in',
            'float': 'floats in',
            'split': 'splits and enters',
            'appear': 'appears instantly'
        },
        'out': {
            'fade': 'fades out of view',
            'fly': 'flies out',
            'wipe': 'wipes out',
            'zoom': 'zooms out',
            'swivel': 'swivels out',
            'bounce': 'bounces out',
            'float': 'floats out',
            'split': 'splits and exits',
            'disappear': 'disappears instantly'
        },
        'emphasis': {
            'color': 'changes color',
            'grow/shrink': 'grows and shrinks',
            'spin': 'spins',
            'pulse': 'pulses',
            'teeter': 'teeters',
            'flash': 'flashes',
            'shimmer': 'shimmers'
        },
        'motion': {
            'path': 'follows a motion path',
            'turn': 'turns',
            'grow': 'grows in size',
            'shrink': 'shrinks in size'
        }
    }
    
    # Get the effect description
    effect_type = animation.get('effect_type', 'appear')
    effect_subtype = animation.get('effect_subtype', '')
    
    if effect_type in effect_descriptions and effect_subtype in effect_descriptions[effect_type]:
        action = effect_descriptions[effect_type][effect_subtype]
    elif effect_type == 'in':
        action = 'enters the slide'
    elif effect_type == 'out':
        action = 'exits the slide'
    elif effect_type == 'emphasis':
        action = 'is emphasized'
    elif effect_type == 'motion':
        action = 'moves'
    else:
        action = 'animates'
    
    # Add direction details
    direction = animation.get('effect_direction', '')
    if direction:
        direction_map = {
            'fromBottom': 'from the bottom',
            'fromTop': 'from the top',
            'fromLeft': 'from the left',
            'fromRight': 'from the right',
            'fromBottomLeft': 'from the bottom-left corner',
            'fromBottomRight': 'from the bottom-right corner',
            'fromTopLeft': 'from the top-left corner',
            'fromTopRight': 'from the top-right corner',
            'horizontal': 'horizontally',
            'vertical': 'vertically',
            'in': 'inward',
            'out': 'outward'
        }
        
        if direction in direction_map:
            action += f" {direction_map[direction]}"
        elif direction.startswith('to_color_'):
            color = direction.replace('to_color_', '#')
            action += f" to the color {color}"
        elif 'scale' in direction:
            # Parse scale values
            import re
            scale_match = re.search(r'scale_x(\d+)_y(\d+)', direction)
            if scale_match:
                x_scale = int(scale_match.group(1)) / 100000
                y_scale = int(scale_match.group(2)) / 100000
                action += f" by {x_scale:.1f}x horizontally and {y_scale:.1f}x vertically"
    
    # Build timing description
    timing_desc = []
    
    # Start condition
    start = animation.get('start_condition', 'on_click')
    if start == 'on_click':
        timing_desc.append("This animation starts when the presenter clicks")
    elif start == 'with_previous':
        timing_desc.append("This animation plays simultaneously with the previous animation")
    elif start == 'after_previous':
        timing_desc.append("This animation starts automatically after the previous animation completes")
    
    # Delay
    delay = animation.get('delay_ms', 0)
    if delay and delay > 0:
        delay_sec = delay / 1000
        if delay_sec >= 1:
            timing_desc.append(f"with a {delay_sec:.1f} second delay")
        else:
            timing_desc.append(f"with a {int(delay)} millisecond delay")
    
    # Duration
    duration = animation.get('duration_ms', 'unknown')
    if duration != 'unknown' and duration != "unknown":
        dur_sec = duration / 1000
        if dur_sec >= 1:
            timing_desc.append(f"taking {dur_sec:.1f} seconds to complete")
        else:
            timing_desc.append(f"taking {int(duration)} milliseconds to complete")
    
    # Repetition
    repeat = animation.get('repeat_count', '1')
    if repeat != '1':
        if repeat == 'indefinite':
            timing_desc.append("repeating continuously")
        else:
            timing_desc.append(f"repeating {repeat} times")
    
    # Auto-reverse
    if animation.get('auto_reverse', False):
        timing_desc.append("then reversing back to its original state")
    
    # Build level (for text animations)
    build_level = animation.get('build_level', '')
    if build_level:
        # Parse paragraph range
        para_match = re.search(r'paragraph_(\d+)-(\d+)', build_level)
        if para_match:
            start_para = int(para_match.group(1))
            end_para = int(para_match.group(2))
            if start_para == end_para:
                timing_desc.append(f"animating paragraph {start_para + 1}")
            else:
                timing_desc.append(f"animating paragraphs {start_para + 1} through {end_para + 1}")
    
    # Combine all parts into a natural description
    full_description = f"{element_desc} {action}."
    
    if timing_desc:
        full_description += " " + ". ".join(timing_desc) + "."
    
    # Add sequence information
    seq_id = animation.get('sequence_id', '')
    effect_id = animation.get('effect_id', '')
    if seq_id and effect_id:
        full_description += f" (Animation sequence {seq_id}, effect {effect_id})"
    
    return full_description

def check_slide_master_animations(pptx_path):
    """
    Check if slide masters and layouts in the presentation contain animations.
    
    Args:
        pptx_path: Path to the PowerPoint file
        
    Returns:
        dict: Dictionary mapping layout indices to boolean indicating if they have animations
    """
    animations_by_layout = {}
    
    try:
        with zipfile.ZipFile(pptx_path) as pptx_zip:
            # Look for slide master files
            master_files = [f for f in pptx_zip.namelist() 
                           if f.startswith('ppt/slideMasters/slideMaster') and f.endswith('.xml')]
            
            for master_file in master_files:
                with pptx_zip.open(master_file) as master_xml:
                    root = ET.parse(master_xml).getroot()
                    if has_animations_in_xml(root):
                        # If master has animations, all its layouts inherit them
                        # Extract number from filename like 'slideMaster1.xml'
                        import re
                        match = re.search(r'slideMaster(\d+)\.xml', master_file)
                        if match:
                            master_idx = match.group(1)
                            animations_by_layout[f'master_{master_idx}'] = True
                        
            # Check slide layouts as well
            layout_files = [f for f in pptx_zip.namelist() 
                           if f.startswith('ppt/slideLayouts/slideLayout') and f.endswith('.xml')]
                           
            for layout_file in layout_files:
                with pptx_zip.open(layout_file) as layout_xml:
                    root = ET.parse(layout_xml).getroot()
                    if has_animations_in_xml(root):
                        # Extract layout index from filename like 'slideLayout12.xml'
                        import re
                        match = re.search(r'slideLayout(\d+)\.xml', layout_file)
                        if match:
                            layout_idx = match.group(1)
                            animations_by_layout[f'layout_{layout_idx}'] = True
    except Exception as e:
        logger.error(f"Error checking slide master animations: {e}", exc_info=True)
        
    return animations_by_layout

def get_slide_layout_info(slide_number, pptx_path):
    """
    Get the layout information for a slide by reading the slide's relationships directly from the zip.
    
    Args:
        slide_number: The slide number (1-based)
        pptx_path: Path to the PowerPoint file
        
    Returns:
        tuple: (layout_index, master_index) or (None, None) if not found
    """
    try:
        import re
        with zipfile.ZipFile(pptx_path) as pptx_zip:
            # Read the slide's relationships file
            slide_rels_path = f'ppt/slides/_rels/slide{slide_number}.xml.rels'
            if slide_rels_path in pptx_zip.namelist():
                with pptx_zip.open(slide_rels_path) as rels_xml:
                    rels_root = ET.parse(rels_xml).getroot()
                    # Look for slideLayout relationship
                    for relationship in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                        target = relationship.get('Target')
                        if target and 'slideLayout' in target:
                            match = re.search(r'slideLayout(\d+)\.xml', target)
                            if match:
                                layout_idx = match.group(1)
                                
                                # Now find which master this layout belongs to
                                layout_rels_path = f'ppt/slideLayouts/_rels/slideLayout{layout_idx}.xml.rels'
                                if layout_rels_path in pptx_zip.namelist():
                                    with pptx_zip.open(layout_rels_path) as layout_rels_xml:
                                        layout_rels_root = ET.parse(layout_rels_xml).getroot()
                                        for rel in layout_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                                            rel_target = rel.get('Target')
                                            if rel_target and 'slideMaster' in rel_target:
                                                master_match = re.search(r'slideMaster(\d+)\.xml', rel_target)
                                                if master_match:
                                                    master_idx = master_match.group(1)
                                                    return (layout_idx, master_idx)
                                return (layout_idx, None)
    except Exception as e:
        logger.debug(f"Could not get layout info for slide {slide_number}: {e}")
    
    return (None, None)

def extract_slide_animations(pptx_path, slide_filter=None):
    """
    Extract animations from all slides in a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        slide_filter (set): Optional set of slide numbers to process
        
    Returns:
        dict: Dictionary containing animation information for all slides
    """
    # Register XML namespaces
    register_namespaces()
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file for animation extraction: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Check which slide masters and layouts contain animations
    animations_by_layout = check_slide_master_animations(pptx_path)
    logger.info(f"Layouts with animations: {animations_by_layout}")
    
    # Dictionary to store animation data
    animation_data = {}
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        # Skip if slide filtering is enabled and this slide is not in the filter
        if slide_filter and i not in slide_filter:
            continue
        # Get slide title
        title = get_slide_title(slide)
        
        # Extract animations directly from slide
        animations = extract_animation_info(slide)
        
        # Also check for animations in the slide XML directly using a simpler method
        has_direct_animations = False
        try:
            with zipfile.ZipFile(pptx_path) as pptx_zip:
                slide_xml_path = f'ppt/slides/slide{i}.xml'
                if slide_xml_path in pptx_zip.namelist():
                    with pptx_zip.open(slide_xml_path) as slide_xml_file:
                        slide_root = ET.parse(slide_xml_file).getroot()
                        if has_animations_in_xml(slide_root):
                            has_direct_animations = True
        except Exception as e:
            logger.debug(f"Could not check direct animations for slide {i}: {e}")
        
        # Get shape information
        shape_info = {}
        for shape in slide.shapes:
            if shape.shape_id:
                shape_type = "Unknown"
                if hasattr(shape, "shape_type"):
                    shape_type = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
                
                shape_text = ""
                if hasattr(shape, "text") and shape.has_text_frame:
                    shape_text = shape.text.strip()
                
                shape_info[str(shape.shape_id)] = {
                    'type': shape_type,
                    'text': shape_text[:100] + ('...' if len(shape_text) > 100 else '')
                }
        
        # Get slide transition
        transition = "None"
        if hasattr(slide, "slide_layout") and hasattr(slide.slide_layout, "transition"):
            transition = str(slide.slide_layout.transition)
        
        # Check for animations in the slide XML directly
        has_slide_animations = len(animations) > 0 or has_direct_animations
        
        # Check if this slide's layout or master has animations
        layout_has_animations = False
        layout_idx, master_idx = get_slide_layout_info(i, pptx_path)
        
        logger.debug(f"Slide {i}: layout_idx={layout_idx}, master_idx={master_idx}")
        
        if layout_idx and f'layout_{layout_idx}' in animations_by_layout:
            layout_has_animations = True
            logger.debug(f"Slide {i} uses layout {layout_idx} which has animations")
        
        # Create animation details with descriptions
        animation_details = []
        for anim in animations:
            anim_detail = anim.copy()
            anim_detail['description'] = create_animation_description(anim, shape_info)
            animation_details.append(anim_detail)
        
        # If slide inherits animations from layout but has no direct animations,
        # try to extract animations from the layout
        if layout_has_animations and len(animations) == 0:
            logger.debug(f"Slide {i} inherits animations from layout {layout_idx}, extracting layout animations")
            try:
                with zipfile.ZipFile(pptx_path) as pptx_zip:
                    layout_path = f'ppt/slideLayouts/slideLayout{layout_idx}.xml'
                    if layout_path in pptx_zip.namelist():
                        with pptx_zip.open(layout_path) as layout_xml:
                            # Parse layout XML and extract animations
                            layout_root = ET.parse(layout_xml).getroot()
                            # Create a mock slide object for the layout
                            class LayoutSlide:
                                def __init__(self, element):
                                    self.element = element
                            
                            layout_slide = LayoutSlide(layout_root)
                            layout_animations = extract_animation_info(layout_slide)
                            
                            # Add layout animations with a note that they're inherited
                            for anim in layout_animations:
                                anim_detail = anim.copy()
                                anim_detail['inherited_from'] = f'layout_{layout_idx}'
                                anim_detail['description'] = f"[Inherited from layout] {create_animation_description(anim, shape_info)}"
                                animation_details.append(anim_detail)
            except Exception as e:
                logger.debug(f"Could not extract animations from layout {layout_idx}: {e}")
        
        # Create animation summary
        animation_summary = ""
        if animation_details:
            # Group animations by sequence
            sequences = {}
            for anim in animation_details:
                seq_id = anim.get('sequence_id', 'unknown')
                if seq_id not in sequences:
                    sequences[seq_id] = []
                sequences[seq_id].append(anim)
            
            # Create narrative summary
            summary_parts = []
            summary_parts.append(f"This slide has {len(animation_details)} animation effects.")
            
            if layout_has_animations and not has_slide_animations:
                summary_parts.append(f"All animations are inherited from the slide layout.")
            elif has_slide_animations and layout_has_animations:
                direct_count = len([a for a in animation_details if 'inherited_from' not in a])
                inherited_count = len([a for a in animation_details if 'inherited_from' in a])
                summary_parts.append(f"{direct_count} animations are directly applied and {inherited_count} are inherited from the layout.")
            
            # Describe the animation flow
            if len(sequences) == 1:
                summary_parts.append("The animations play in a single sequence.")
            else:
                summary_parts.append(f"The animations are organized in {len(sequences)} sequences.")
            
            animation_summary = " ".join(summary_parts)
        else:
            animation_summary = "This slide has no animations."
        
        # Add slide information to the dictionary
        animation_data[f"slide_{i}"] = {
            'slide_number': i,
            'title': title,
            'animations': animations,
            'animation_details': animation_details,
            'animation_summary': animation_summary,
            'shapes': shape_info,
            'transition': transition,
            'animation_count': len(animation_details),
            'has_animations': has_slide_animations or layout_has_animations,
            'layout_animations': layout_has_animations,
            'direct_animations': has_slide_animations
        }
        
        logger.info(f"Processed slide {i}: {title[:50]}{'...' if len(title) > 50 else ''} - Direct animations: {len(animations)}, Layout animations: {layout_has_animations}")
    
    return animation_data
