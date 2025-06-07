#!/usr/bin/env python3
"""
PowerPoint Notes Extractor
-------------------------
This script extracts notes from a PowerPoint file and saves them to a JSON file.
"""

import os
import sys
import json
import argparse
import logging
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Define XML namespaces used in PPTX files
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006'
}

def register_namespaces():
    """Register XML namespaces for parsing."""
    for prefix, uri in NAMESPACES.items():
        ET.register_namespace(prefix, uri)

def extract_notes_from_xml(pptx_path):
    """
    Extract notes directly from the PPTX XML structure.
    
    This function parses the notes XML files in a PowerPoint presentation to extract
    notes text. It looks for text within the body placeholder in notes slides.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        dict: Dictionary mapping slide numbers to notes text
    """
    notes_by_slide = {}
    
    try:
        with zipfile.ZipFile(pptx_path) as pptx_zip:
            # Get list of all notes files
            notes_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/notesSlides/notesSlide') and f.endswith('.xml')]
            
            if not notes_files:
                logger.info("No notes files found in the PowerPoint file.")
                return notes_by_slide
            
            logger.info(f"Found {len(notes_files)} notes files in the PowerPoint file.")
            
            # Process each notes file
            for notes_file in notes_files:
                try:
                    # Extract slide number from notes file name
                    slide_num = int(notes_file.split('notesSlide')[1].split('.')[0])
                    
                    # Extract notes text from XML
                    with pptx_zip.open(notes_file) as notes_xml:
                        xml_content = notes_xml.read().decode('utf-8')
                        logger.debug(f"Processing notes file: {notes_file}")
                        
                        # Parse XML content
                        root = ET.fromstring(xml_content)
                        
                        # Find all shapes in the notes slide
                        shapes = root.findall('.//p:sp', NAMESPACES)
                        
                        # Look for the shape with the notes content (has placeholder type="body")
                        notes_text = ""
                        for shape in shapes:
                            # Check if this is the notes placeholder
                            ph_elem = shape.find('.//p:nvPr/p:ph[@type="body"]', NAMESPACES)
                            if ph_elem is not None:
                                # This is the notes placeholder, extract text
                                tx_body = shape.find('.//p:txBody', NAMESPACES)
                                if tx_body is not None:
                                    # Extract all paragraphs
                                    paragraphs = []
                                    
                                    # Find all paragraph elements
                                    p_elems = tx_body.findall('.//a:p', NAMESPACES)
                                    for p in p_elems:
                                        # Find all text runs in this paragraph
                                        r_elems = p.findall('.//a:r', NAMESPACES)
                                        para_text = []
                                        
                                        # If there are no text runs, check for direct text elements
                                        if not r_elems:
                                            t_elems = p.findall('.//a:t', NAMESPACES)
                                            for t in t_elems:
                                                if t.text and t.text.strip():
                                                    para_text.append(t.text)
                                        else:
                                            # Extract text from each run
                                            for r in r_elems:
                                                t_elems = r.findall('.//a:t', NAMESPACES)
                                                for t in t_elems:
                                                    if t.text and t.text.strip():
                                                        para_text.append(t.text)
                                        
                                        # Add paragraph text if not empty
                                        if para_text:
                                            paragraphs.append(' '.join(para_text))
                                    
                                    # Combine paragraphs into notes text
                                    if paragraphs:
                                        notes_text = '\n'.join(paragraphs)
                                        notes_by_slide[slide_num] = notes_text
                                        logger.info(f"Found notes for slide {slide_num}: {notes_text[:50]}..." if len(notes_text) > 50 else f"Found notes for slide {slide_num}: {notes_text}")
                                        break  # Found the notes, no need to check other shapes
                        
                        # Log if no notes content was found
                        if slide_num not in notes_by_slide:
                            logger.debug(f"No notes content found in {notes_file}")
                
                except Exception as e:
                    logger.error(f"Error processing notes file {notes_file}: {e}", exc_info=True)
    
    except Exception as e:
        logger.error(f"Error extracting notes from XML: {e}", exc_info=True)
    
    return notes_by_slide

def extract_slide_notes(pptx_path, output_file):
    """
    Extract notes from all slides in a PowerPoint file and save to JSON.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_file (str): Path to save the JSON output
    
    Returns:
        dict: Dictionary containing notes for all slides
    """
    # Create output directory if it doesn't exist
    output_path = Path(output_file).parent
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Register XML namespaces
    register_namespaces()
    
    # Load the presentation
    logger.info(f"Opening PowerPoint file: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to open PowerPoint file: {e}")
        return None
    
    # Extract notes using direct XML parsing
    logger.info("Extracting notes from XML structure...")
    xml_notes = extract_notes_from_xml(pptx_path)
    
    # Dictionary to store notes information
    notes_data = {}
    
    # Process each slide
    logger.info(f"Found {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title
        title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                if shape.text.strip():
                    title = shape.text.strip().replace('\n', ' ')
                    break
        
        # Extract notes text using python-pptx
        pptx_notes = ""
        if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
            if hasattr(slide.notes_slide, 'notes_text_frame'):
                pptx_notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Get notes from XML parsing
        xml_note = xml_notes.get(i, "")
        
        # Use the best available notes (prefer XML parsing if it found notes)
        notes_text = xml_note if xml_note else pptx_notes
        
        logger.info(f"Processing slide {i}: {title} - Notes: {'Yes' if notes_text else 'No'}")
        
        # Add slide information to the dictionary
        notes_data[f"slide_{i}"] = {
            'slide_number': i,
            'title': title,
            'notes': notes_text
        }
    
    # Save notes information to JSON file
    logger.info(f"Saving notes information to {output_file}")
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(notes_data, f, indent=2)
        logger.info(f"Successfully saved notes information to {output_file}")
    except Exception as e:
        logger.error(f"Failed to save notes information: {e}")
        return None
    
    return notes_data

def main():
    parser = argparse.ArgumentParser(description='Extract notes from a PowerPoint file.')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    parser.add_argument('--output', '-o', default='./slide_notes.json', help='Output JSON file')
    
    args = parser.parse_args()
    
    # Check if the PowerPoint file exists
    if not os.path.isfile(args.pptx_file):
        logger.error(f"PowerPoint file not found: {args.pptx_file}")
        sys.exit(1)
    
    result = extract_slide_notes(args.pptx_file, args.output)
    if result is None:
        sys.exit(1)

if __name__ == '__main__':
    main()
